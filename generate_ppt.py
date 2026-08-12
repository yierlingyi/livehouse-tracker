#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""生成《乐队演出查看小程序》产品介绍 PPT。

运行:
    python generate_ppt.py

输出:
    docs/presentation.pptx

说明:
    面向非技术用户的 6 页产品介绍，深色背景 + 活力橙/金主色调。
    核心入口为 create_pptx() 函数。
"""

from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# 配色：深蓝黑背景 + 活力珊瑚红/金色点缀
# ---------------------------------------------------------------------------
BG = RGBColor(0x1A, 0x1A, 0x2E)      # 深蓝黑
CARD = RGBColor(0x24, 0x24, 0x48)    # 卡片底（略亮）
CARD_HI = RGBColor(0x2E, 0x2E, 0x5A) # 高亮卡片底
ACCENT = RGBColor(0xE9, 0x45, 0x60)  # 活力珊瑚红
GOLD = RGBColor(0xF0, 0xA5, 0x00)    # 金色
TEXT = RGBColor(0xFF, 0xFF, 0xFF)    # 主文字（白）
MUTED = RGBColor(0xB8, 0xB8, 0xD6)   # 次要文字（淡紫灰）

FONT = "微软雅黑"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

MARGIN = Inches(0.7)                 # 左右页边距
GAP = Inches(0.3)                    # 卡片间距


# ---------------------------------------------------------------------------
# 基础工具函数
# ---------------------------------------------------------------------------
def _set_ea(run, name: str):
    """设置中文字体（East Asian typeface）。"""
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", name)


def _fill_bg(slide):
    """整页铺满深色背景。"""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = BG
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _card(slide, x, y, w, h, fill=CARD, line=None, radius=0.07):
    """圆角卡片。"""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    return shp


def _txt(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    """添加文本框。

    paras: 段落列表，每项为 dict：
        {text, size, bold, color, align, line_spacing, space_before, space_after}
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for p in paras:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = p.get("align", PP_ALIGN.LEFT)
        para.line_spacing = p.get("line_spacing", 1.0)
        if p.get("space_before") is not None:
            para.space_before = Pt(p["space_before"])
        if p.get("space_after") is not None:
            para.space_after = Pt(p["space_after"])
        run = para.add_run()
        run.text = p["text"]
        f = run.font
        f.name = FONT
        _set_ea(run, FONT)
        f.size = Pt(p.get("size", 14))
        f.bold = p.get("bold", False)
        f.color.rgb = p.get("color", TEXT)
    return tb


def _circle(slide, x, y, d, fill, line=None):
    """圆形。"""
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def _title(slide, text, y=Inches(0.55), size=34, x=None):
    """统一页面标题。"""
    if x is None:
        x = MARGIN
    _txt(
        slide, x, y, EMU_W - x - MARGIN, Inches(0.9),
        [{"text": text, "size": size, "bold": True, "color": TEXT}],
    )


def _emoji_cell(slide, x, y, w, h, emoji, title, desc, emoji_size=24):
    """核心功能卡片：左侧 emoji，右侧标题 + 描述。"""
    _card(slide, x, y, w, h)
    # 左侧 emoji
    _txt(
        slide, x + Inches(0.28), y, Inches(0.9), h,
        [{"text": emoji, "size": emoji_size, "color": TEXT}],
        anchor=MSO_ANCHOR.MIDDLE, wrap=False,
    )
    # 右侧文字
    _txt(
        slide, x + Inches(1.15), y + Inches(0.18), w - Inches(1.45), h - Inches(0.36),
        [
            {"text": title, "size": 17, "bold": True, "color": TEXT,
             "space_after": 3},
            {"text": desc, "size": 12.5, "color": MUTED, "line_spacing": 1.1},
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )


# ---------------------------------------------------------------------------
# 每页构建
# ---------------------------------------------------------------------------
def _slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
    _fill_bg(s)

    # 顶部小标签
    _txt(
        s, MARGIN, Inches(0.75), EMU_W - 2 * MARGIN, Inches(0.4),
        [{"text": "产 品 介 绍  ·  乐 队 演 出 查 看 小 程 序",
          "size": 14, "bold": True, "color": GOLD, "align": PP_ALIGN.CENTER}],
    )

    # 装饰音符圆
    _circle(s, Inches(6.32), Inches(1.55), Inches(0.7), ACCENT)
    _txt(
        s, Inches(6.32), Inches(1.72), Inches(0.7), Inches(0.5),
        [{"text": "♪", "size": 24, "bold": True, "color": TEXT,
          "align": PP_ALIGN.CENTER}],
    )

    # 主标题
    _txt(
        s, MARGIN, Inches(2.55), EMU_W - 2 * MARGIN, Inches(1.0),
        [{"text": "乐队演出查看小程序", "size": 44, "bold": True,
          "color": TEXT, "align": PP_ALIGN.CENTER}],
    )
    # 副标题
    _txt(
        s, MARGIN, Inches(3.6), EMU_W - 2 * MARGIN, Inches(0.55),
        [{"text": "随时随地，发现身边的现场音乐", "size": 21,
          "color": GOLD, "align": PP_ALIGN.CENTER}],
    )
    # 一句话介绍
    _txt(
        s, MARGIN, Inches(4.35), EMU_W - 2 * MARGIN, Inches(0.6),
        [{"text": "从 Livehouse 到音乐节，一站式掌握身边每一场正在售票的现场演出。",
          "size": 16, "color": MUTED, "align": PP_ALIGN.CENTER}],
    )

    # 受众标签
    labels = ["🎧  游客", "🎸  乐队成员", "🏟️  Livehouse 老板"]
    label_w = Inches(3.0)
    total = label_w * 3 + GAP * 2
    x0 = (EMU_W - total) // 2
    y = Inches(5.45)
    for i, lb in enumerate(labels):
        lx = x0 + i * (label_w + GAP)
        chip = _card(s, lx, y, label_w, Inches(0.62), fill=CARD,
                     line=RGBColor(0x3A, 0x3A, 0x6A), radius=0.5)
        tf = chip.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = lb
        f = run.font
        f.name = FONT
        _set_ea(run, FONT)
        f.size = Pt(14)
        f.color.rgb = TEXT


def _slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(s)
    _title(s, "解决了什么问题？")

    pains = [
        ("痛点 1", "想听现场音乐，却不知道去哪找"),
        ("痛点 2", "演出信息分散在各个社交平台"),
        ("痛点 3", "到了现场，才发现票早已售罄"),
    ]
    card_w = Inches(3.9)
    x0 = (EMU_W - (card_w * 3 + GAP * 2)) // 2
    y = Inches(1.95)
    card_h = Inches(2.5)
    for i, (tag, text) in enumerate(pains):
        cx = x0 + i * (card_w + GAP)
        _card(s, cx, y, card_w, card_h)
        # 红色感叹圆 + 标签
        _circle(s, cx + Inches(0.3), y + Inches(0.32), Inches(0.5), ACCENT)
        _txt(
            s, cx + Inches(0.3), y + Inches(0.42), Inches(0.5), Inches(0.4),
            [{"text": "!", "size": 20, "bold": True, "color": TEXT,
              "align": PP_ALIGN.CENTER}],
        )
        _txt(
            s, cx + Inches(0.95), y + Inches(0.36), card_w - Inches(1.2),
            Inches(0.45),
            [{"text": tag, "size": 15, "bold": True, "color": ACCENT}],
        )
        _txt(
            s, cx + Inches(0.3), y + Inches(1.15), card_w - Inches(0.6),
            Inches(1.1),
            [{"text": text, "size": 18, "bold": True, "color": TEXT,
              "line_spacing": 1.15}],
        )

    # 解决方案横幅
    by = Inches(5.0)
    bw = Inches(12.0)
    bx = (EMU_W - bw) // 2
    bh = Inches(1.7)
    _card(s, bx, by, bw, bh, fill=CARD_HI, line=GOLD)
    _circle(s, bx + Inches(0.4), by + Inches(0.5), Inches(0.7), GOLD)
    _txt(
        s, bx + Inches(0.4), by + Inches(0.66), Inches(0.7), Inches(0.5),
        [{"text": "✓", "size": 24, "bold": True, "color": BG,
          "align": PP_ALIGN.CENTER}],
    )
    _txt(
        s, bx + Inches(1.4), by + Inches(0.32), bw - Inches(1.8), bh - Inches(0.5),
        [
            {"text": "我们的解决方案", "size": 15, "bold": True,
             "color": GOLD, "space_after": 5},
            {"text": "一站式查看所有正在售票的演出",
             "size": 24, "bold": True, "color": TEXT},
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )


def _slide_features(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(s)
    _title(s, "核心功能")

    feats = [
        ("📅", "按日期浏览演出", "今天、本周、周末，想看的演出一目了然"),
        ("🏙️", "切换城市查看演出", "一键切换城市，发现不同地方的演出"),
        ("🎸", "查看演出详情", "乐队阵容、票价、购票链接一次看全"),
        ("⚡", "打开即看，无需登录", "点开链接就能用，不打扰不等待"),
        ("📶", "弱网也能用", "有缓存，不怕没信号"),
    ]

    col_w = Inches(5.9)
    col1 = (EMU_W - (col_w * 2 + GAP)) // 2
    col2 = col1 + col_w + GAP
    row_h = Inches(1.72)
    r1 = Inches(1.9)
    r2 = Inches(3.82)

    _emoji_cell(s, col1, r1, col_w, row_h, *feats[0])
    _emoji_cell(s, col2, r1, col_w, row_h, *feats[1])
    _emoji_cell(s, col1, r2, col_w, row_h, *feats[2])
    _emoji_cell(s, col2, r2, col_w, row_h, *feats[3])

    # 第 5 项：居中宽卡片
    w5 = Inches(9.6)
    x5 = (EMU_W - w5) // 2
    y5 = Inches(5.78)
    _emoji_cell(s, x5, y5, w5, Inches(1.3), *feats[4])


def _slide_scenarios(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(s)
    _title(s, "使用场景")

    scenarios = [
        ("🧳", "游客", ACCENT, [
            "周末想找演出 → 打开小程序",
            "选城市 → 浏览演出列表",
            "查看详情 → 直接买票",
        ]),
        ("🎸", "乐队成员", GOLD, [
            "分享小程序链接给粉丝",
            "粉丝看到详细演出信息",
            "更多人买票，演出更精彩",
        ]),
        ("🏟️", "Livehouse 老板", RGBColor(0x6C, 0x8C, 0xF5), [
            "演出信息集中展示",
            "提高演出曝光度",
            "吸引更多观众到场",
        ]),
    ]

    card_w = Inches(3.9)
    x0 = (EMU_W - (card_w * 3 + GAP * 2)) // 2
    y = Inches(1.9)
    card_h = Inches(4.9)

    for i, (emoji, name, color, steps) in enumerate(scenarios):
        cx = x0 + i * (card_w + GAP)
        _card(s, cx, y, card_w, card_h)

        # 顶部：emoji 圆 + 受众名
        _circle(s, cx + Inches(0.32), y + Inches(0.35), Inches(0.7), color)
        _txt(
            s, cx + Inches(0.32), y + Inches(0.52), Inches(0.7), Inches(0.45),
            [{"text": emoji, "size": 22, "align": PP_ALIGN.CENTER}],
        )
        _txt(
            s, cx + Inches(1.2), y + Inches(0.42), card_w - Inches(1.5),
            Inches(0.6),
            [{"text": name, "size": 18, "bold": True, "color": TEXT}],
        )

        # 步骤
        paras = []
        for j, st in enumerate(steps):
            paras.append({
                "text": f"▸ {st}", "size": 14.5, "color": TEXT,
                "space_before": 14, "line_spacing": 1.12,
            })
        _txt(
            s, cx + Inches(0.34), y + Inches(1.4), card_w - Inches(0.65),
            Inches(3.3), paras,
        )


def _slide_easy(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(s)
    _title(s, "简单好用")

    items = [
        ("无需注册，打开即用", "不用填表，扫一扫或点开链接就能看"),
        ("数据自动更新，下拉刷新即可", "最新排期实时同步，刷新一下就有"),
        ("离线也能查看已加载的演出", "提前加载过的内容，没信号也能看"),
        ("界面简洁清晰，一看就懂", "没有复杂操作，谁都能轻松上手"),
    ]

    col_w = Inches(5.9)
    col1 = (EMU_W - (col_w * 2 + GAP)) // 2
    col2 = col1 + col_w + GAP
    row_h = Inches(2.35)
    r1 = Inches(1.9)
    r2 = Inches(4.45)

    for i, (title, desc) in enumerate(items):
        cx = col1 if i % 2 == 0 else col2
        cy = r1 if i < 2 else r2
        _card(s, cx, cy, col_w, row_h)
        _circle(s, cx + Inches(0.32), cy + Inches(0.32), Inches(0.62), GOLD)
        _txt(
            s, cx + Inches(0.32), cy + Inches(0.48), Inches(0.62), Inches(0.4),
            [{"text": "✓", "size": 20, "bold": True, "color": BG,
              "align": PP_ALIGN.CENTER}],
        )
        _txt(
            s, cx + Inches(1.2), cy + Inches(0.36), col_w - Inches(1.5),
            Inches(1.7),
            [
                {"text": title, "size": 18, "bold": True, "color": TEXT,
                 "space_after": 8},
                {"text": desc, "size": 13.5, "color": MUTED,
                 "line_spacing": 1.15},
            ],
            anchor=MSO_ANCHOR.TOP,
        )


def _slide_closing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(s)

    # 装饰音符圆
    _circle(s, Inches(6.32), Inches(1.15), Inches(0.7), ACCENT)
    _txt(
        s, Inches(6.32), Inches(1.32), Inches(0.7), Inches(0.5),
        [{"text": "♪", "size": 24, "bold": True, "color": TEXT,
          "align": PP_ALIGN.CENTER}],
    )

    # 金句
    _txt(
        s, MARGIN, Inches(2.15), EMU_W - 2 * MARGIN, Inches(0.9),
        [{"text": "让每一场演出都被看见", "size": 40, "bold": True,
          "color": TEXT, "align": PP_ALIGN.CENTER}],
    )
    _txt(
        s, MARGIN, Inches(3.2), EMU_W - 2 * MARGIN, Inches(0.5),
        [{"text": "乐队演出查看小程序 · 产品介绍", "size": 16,
          "color": MUTED, "align": PP_ALIGN.CENTER}],
    )

    # 联系方式卡片（留白供填写）
    cw = Inches(8.4)
    cx = (EMU_W - cw) // 2
    cy = Inches(4.15)
    ch = Inches(2.5)
    _card(s, cx, cy, cw, ch, fill=CARD, line=RGBColor(0x3A, 0x3A, 0x6A))
    _txt(
        s, cx + Inches(0.5), cy + Inches(0.35), cw - Inches(1.0), Inches(0.5),
        [{"text": "联系我们", "size": 17, "bold": True, "color": GOLD}],
    )
    contact_paras = [
        {"text": "微信公众号：________________", "size": 15,
         "color": TEXT, "space_before": 12},
        {"text": "电子邮箱：________________", "size": 15,
         "color": TEXT, "space_before": 10},
        {"text": "联系电话：________________", "size": 15,
         "color": TEXT, "space_before": 10},
    ]
    _txt(
        s, cx + Inches(0.5), cy + Inches(0.85), cw - Inches(1.0), Inches(1.5),
        contact_paras,
    )


# ---------------------------------------------------------------------------
# 入口
# ---------------------------------------------------------------------------
def create_pptx(out_path: str | os.PathLike = os.path.join("docs", "presentation.pptx")) -> str:
    """生成产品介绍 PPT，返回输出文件路径。"""
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H

    _slide_cover(prs)       # 1 封面
    _slide_problem(prs)     # 2 解决了什么问题
    _slide_features(prs)    # 3 核心功能
    _slide_scenarios(prs)   # 4 使用场景
    _slide_easy(prs)        # 5 简单好用
    _slide_closing(prs)     # 6 让每一场演出都被看见

    out_path = os.path.abspath(out_path)
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    path = create_pptx()
    print(f"已生成: {path}")
